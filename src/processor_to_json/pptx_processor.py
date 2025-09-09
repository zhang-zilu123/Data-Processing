# PPTX演示文稿工厂信息处理模块(具有固定格式的信息表)
# 功能：解析PPT演示文稿中的工厂信息，提取文本内容和微信二维码，转换为标准JSON格式
# 支持多页处理、文本智能识别、AI模型解析、微信二维码提取、数据标准化转换

import re
import logging
from pptx import Presentation
from typing import List

import sys
import os
# 添加项目根目录到路径
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))
from setting.config import *  # 导入配置模块
from src.processor_to_json.processor_rely.model_remark_pptx_info import extract_remark_info
from src.processor_to_json.processor_rely.outmodel_results_validator import validate_and_get_result
from src.utils.clean_factory_name import clean_factory_name
from src.utils.save_result_to_json import make_vendor_folder,save_result_to_vendor_folder
from src.utils.SaveImg_wechat_qr import extract_images_from_pptx

# 日志配置
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


#---------------------- 文本识别和分类工具函数 --------------------------------

# --- 日期文本格式识别函数 ---
def is_date_text(text:str) -> bool:
    """
    检查文本内容是否符合日期格式模式
    
    识别功能：
    通过config中DATE_PATTERNS配置的日期正则表达式模式，判断输入文本是否为日期信息。
    用于在文档解析过程中识别和分类日期相关的内容。
    
    识别流程：
    1. 基础文本有效性检查（长度和非空）
    2. 清理文本格式（去除首尾空白）
    3. 遍历配置的日期模式进行匹配
    4. 返回匹配结果
    
    参数：
        text (str): 待检查的文本内容
        
    返回：
        bool: True表示文本符合日期格式，False表示不符合

    """

    # 步骤1：基础有效性检查
    if not text or len(text) > 50: # 简单的长度过滤，避免检查过长的文本
        return False
        
    # 步骤2：清理文本格式
    cleaned_text = text.strip()
    
    # 步骤3：遍历配置的日期模式进行匹配
    for pattern in DATE_PATTERNS:
        if re.search(pattern, cleaned_text):
            return True
    return False



# --- 联系人和电话号码识别函数 ---
def find_contact_line(text:str) -> str:
    """
    识别文本中的联系方式信息（电话号码、职位、姓名）
    
    识别功能：
    使用正则表达式模式识别文本中的电话号码，支持多种格式的电话号码
    （包括带空格、短横线分隔的11位数字）。用于自动提取联系方式行。
    
    识别流程：
    1. 编译电话号码识别的正则表达式
    2. 在输入文本中搜索匹配模式
    3. 返回匹配结果或空值
    
    参数：
        text (str): 待搜索的文本内容
        
    返回：
        str: 如果找到电话号码则返回原始文本，否则返回None

    """
    # 匹配手机号的正则模式（11 位数字，允许有空格或短横线）
    phone_pattern = re.compile(r"(?<!\d)(?:\d[\s\-]*){11}(?!\d)")
    
    # 放宽的职位关键词匹配模式（支持空格和多种职位）
    position_keywords = ['经理', '负责人', '主管', '总监', '主任', '部长', '董事', '总','总经理','助理','秘书','代表','销售','业务']
    position_pattern = re.compile(r'.*(' + '|'.join(position_keywords) + r').*')

    # 匹配2-4个中文字的模式
    chinese_name_pattern = re.compile(r'^[\u4e00-\u9fa5]{2,4}$')

    # 检查是否匹配手机号、职位关键词或2-4个中文字
    if phone_pattern.search(text) or position_pattern.search(text) or chinese_name_pattern.search(text):
        return text
    else:
        return None

# --- 厂商名称格式识别函数 ---
def is_vendor_name(line:str) -> bool:
    """
    判断文本行是否为厂商名称格式
    
    识别功能：
    通过检查文本特征（无冒号且包含企业关键词），判断当前行
    是否为厂商名称信息。用于自动识别和分类厂商名称字段。
    
    识别流程：
    1. 检查文本中是否包含冒号
    2. 检查是否包含企业类型关键词
    3. 综合判断返回结果
    
    参数：
        line (str): 待判断的文本行
        
    返回：
        bool: True表示可能是厂商名称，False表示不是
        

    """
    # 步骤1：检查是否包含冒号（有冒号的通常是标签：值格式）
    if ':' in line or ':' in line:
        return False
        
    # 步骤2：检查是否包含企业类型关键词
    if any(kw in line for kw in ['公司', '厂', '集团', '有限', '企业']):
        return True
    return False


#---------------------- PPTX文本提取和处理模块 --------------------------------

# --- 单页幻灯片文本提取函数 ---
def extract_text_from_slide(slide) -> List[str]:
    """
    从单个PPT幻灯片中提取所有文本内容
    
    提取流程：
    1. 遍历幻灯片中的所有形状
    2. 检查形状是否包含文本框
    3. 提取文本框中的所有段落文本
    4. 过滤空行并返回文本行列表
    
    参数：
        slide: PPT幻灯片对象
            
    返回：
        List[str]: 包含该幻灯片所有文本行的列表
    """
    
    # 收集所有文本元素及其位置
    elements = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            # 处理多行文本
            full_text = ""
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    full_text += run.text
                full_text += "\n"  # 保留段落分隔
            full_text = full_text.strip()
            
            if full_text:
                # 记录元素位置和文本
                elements.append({
                    'top': shape.top,
                    'height': shape.height,
                    'text': full_text
                })
    
    # 按垂直位置排序（从上到下）
    elements.sort(key=lambda x: x['top'])
    
    # 重组文本块（处理多行内容）
    ordered_texts = []
    for el in elements:
        # 分割多行文本
        lines = el['text'].split('\n')
        for line in lines:
            if line.strip():  # 忽略空行
                ordered_texts.append(line.strip())
    
    # 特殊处理：确保日期在最后
    date_pattern = r'\d{4}/\d{2}/\d{2}'
    for i, text in enumerate(ordered_texts):
        if re.match(date_pattern, text):
            # 移除并添加到末尾
            date_text = ordered_texts.pop(i)
            ordered_texts.append(date_text)
            break
    
    return ordered_texts


# --- PPTX文件全页文本提取函数 ---
def extract_text_from_pptx(file_path: str) -> List[List[str]]:
    """
    从PPTX文件中逐页提取所有文本内容
    
    提取流程：
    1. 打开PPTX文件
    2. 遍历所有幻灯片
    3. 逐页提取文本内容
    4. 返回每页的文本行列表
    
    参数：
        file_path: PPTX文件路径
        
    返回：
        List[List[str]]: 包含所有幻灯片文本行的二维列表，每个元素是一页的文本行列表
    """
    try:
        logging.info(f"开始解析PPTX文件: {file_path}")
        prs = Presentation(file_path)
        all_slides_text = []
        
        # 遍历所有幻灯片
        for slide_num, slide in enumerate(prs.slides, 1):
            
            # 提取当前幻灯片的文本
            slide_text = extract_text_from_slide(slide)
            
            if slide_text:
                all_slides_text.append(slide_text)
            else:
                logging.error(f"第 {slide_num} 页没有提取到文本")
        
        return all_slides_text
        
    except Exception as e:
        logging.error(f"解析PPTX文件失败: {str(e)}")
        raise


#---------------------- 文本解析和JSON转换模块 --------------------------------

# --- 文本行转JSON格式函数 ---
def extract_text_to_json(text_lines: list) -> dict:
    """
    使用智能匹配逻辑将文本行转换为JSON格式
    
    转换逻辑：
    当识别到关键词时，将后续数据作为该字段的值，直到遇到下一个关键词。
    支持厂商名称、联系方式、日期等特殊字段的自动识别。
    
    转换流程：
    1. 初始化字段结果字典
    2. 前3行进行厂商名称匹配
    3. 接下来进行联系方式匹配
    4. 关键词匹配和字段填充
    5. 日期识别和处理
    
    参数：
        text_lines (list): 文本行列表
        
    返回：
        dict: 转换后的JSON格式字典
    """
    try:
        # 步骤1：初始化字段结果字典，使用配置的标准字段
        result = {field: '' for field in JSON_FORMAT.keys()}
        
        i = 0
        n = len(text_lines)
        current_field = None  # 当前正在填充的字段
        
        # 1. 厂商名称匹配（前3行）
        while i < min(3, n):
            # 预处理
            line = text_lines[i].strip()
            text_lines[i] = line

            # 默认第一行为厂商名称
            if i == 0:
                result['厂商名称'] = line
                i += 1
                continue

            # 厂商名称匹配
            if is_vendor_name(line):
                if result['厂商名称'] == '':
                    result['厂商名称'] = line
                else:
                    result['厂商名称'] += '/' + line
                i += 1
            else:
                break
        
        # 2. 联系方式匹配（接下来3行）
        while i < min(8, n):
            line = text_lines[i].strip()
            if find_contact_line(line):
                if result['联系方式'] == '':
                    result['联系方式'] = line
                else:
                    result['联系方式'] += '\n' + line
                i += 1
            else:
                break
        

        # 4. 关键词匹配
        while i < n-1:
            line = text_lines[i].strip()
            if not line:  # 跳过空行
                i += 1
                continue
            
            # 检查当前行是否包含关键词
            matched_field = None
            for field_name, keywords in TEXT_LABELS_pptx.items():
                if any(keyword in line for keyword in keywords):
                    matched_field = field_name
                    break
            
            if matched_field:
                # 找到关键词，切换到新字段
                current_field = matched_field
                i += 1  # 移动到下一行（关键词行不作为内容）
                continue
            else:
                # 当前行不是关键词
                if current_field:
                    # 有当前字段，将内容填入该字段
                    if result[current_field] == '':
                        result[current_field] = line
                    else:
                        result[current_field] += '\n' + line
                else:
                    # 没有当前字段，填入备注
                    if result['备注'] == '':
                        result['备注'] = line
                    else:
                        result['备注'] += '\n' + line
                
                i += 1
            # 3. 日期匹配
            if i == n-1:
                if is_date_text(text_lines[i]):
                    result['日期'] = text_lines[i]
                    i += 1
                break

        
        return result
        
    except Exception as e:
        logging.error(f"转换为JSON格式失败: {str(e)}")
        return {}

# --- 备注信息AI解析函数 ---
def extract_info_remarks(text_lines: dict) -> dict:
    """
    使用AI模型提取备注信息并更新字段
    
    解析功能：
    调用AI模型对备注字段进行深度解析，提取其中包含的结构化信息，
    并将解析结果更新到对应的字段中。
    
    解析流程：
    1. 获取备注字段内容
    2. 循环3次调用AI模型进行解析
    3. 验证模型输出的一致性
    4. 更新解析结果到原字段
    5. 返回更新后的数据
    
    参数：
        text_lines (dict): 包含备注信息的字典
        
    返回：
        dict: 更新后的字典，失败时返回原始数据
    """
    try:
        # 获取备注字段内容
        remarks = text_lines['备注']
        
        if not remarks:
            return text_lines
        
        # 循环3次，每次调用2次模型
        for attempt in range(3):
            logging.info(f"第{attempt + 1}次调用模型解析备注字段...")
            
            # 第一次调用模型
            remark_result_1 = extract_remark_info(remarks)
            
            # 第二次调用模型
            remark_result_2 = extract_remark_info(remarks)
            
            # 检查API调用是否成功
            if remark_result_1 is None or remark_result_1 == '':
                logging.error(f"第{attempt + 1}次尝试 - 第一次API调用失败，返回空结果")
                continue
                
            if remark_result_2 is None or remark_result_2 == '':
                logging.error(f"第{attempt + 1}次尝试 - 第二次API调用失败，返回空结果")
                continue
            
            # 使用通用工具验证并获取结果
            remark_data = validate_and_get_result([remark_result_1, remark_result_2])
            
            if remark_data is not None:
                # 更新原字段
                updated_result = text_lines.copy()
                
                # 遍历模型返回的数据，直接匹配字段名称
                for field, value in remark_data.items():
                    if not value or not str(value).strip():
                        continue
                        
                    if field in updated_result:
                        # 如果字段存在且不是备注字段，追加内容
                        if field != '备注':
                            if updated_result[field]:
                                updated_result[field] += '\n' + str(value)
                            else:
                                updated_result[field] = str(value)
                        else:
                            # 备注字段直接覆盖
                            updated_result[field] = str(value)
                    else:
                        # 如果字段不存在，直接添加
                        updated_result[field] = str(value)
                        
                    logging.info(f"字段更新: {field} -> {value}")
                
                return updated_result
        
        # 如果3次循环都没有找到一致的结果，报错并返回原始数据
        logging.error("3次尝试都失败，模型调用结果不一致，跳过备注处理")
        return text_lines
    
    except Exception as e:
        logging.error(f"处理备注信息时出错: {e}")
        return text_lines


#---------------------- PPTX文件综合处理主模块 --------------------------------
        
# --- PPTX文件主处理函数 ---
def process_pptx_file(file_path: str,output_directory:str) -> bool:
    """
    PPTX文件完整处理流程，逐页处理并保存结果
    
    处理功能：
    处理PPTX演示文稿文件，逐页提取工厂信息，转换为标准JSON格式，
    提取微信二维码图片，并按厂商分类保存到指定目录。
    
    处理流程：
    1. 逐页提取PPTX文件中的文本
    2. 对每页文本进行JSON转换
    3. 处理每页的备注信息
    4. 创建厂商文件夹并保存结果
    5. 提取并保存微信二维码图片
    
    参数：
        file_path (str): PPTX文件路径
        output_directory (str): 输出目录路径
        
    返回：
        bool: 处理成功返回True，失败返回False
    """
    try:
        # 步骤1：逐页提取文本
        all_slides_text = extract_text_from_pptx(file_path)
        
        # 步骤2：逐页处理
        results = []
        for slide_num, slide_text in enumerate(all_slides_text, 1):
            
            # 转换为JSON格式
            json_result = extract_text_to_json(slide_text)
            
            if json_result:
                # 处理备注信息
                final_result = extract_info_remarks(json_result)
                results.append(final_result)
                logging.info(f"第 {slide_num} 页处理完成")
            else:
                logging.warning(f"第 {slide_num} 页转换失败")
        
        logging.info(f"PPTX文件处理完成，共处理 {len(results)} 页")
       
        # 步骤3：保存处理结果
        if results and isinstance(results, list):
            success_count = 0
            for i, result in enumerate(results):
                try:
                    # 检查result是否为字典且包含厂商名称
                    if not isinstance(result, dict):
                        logging.warning(f"第{i+1}页结果不是字典格式，跳过")
                        continue
                    
                    # 获取厂商名称
                    vendor_name = result.get('厂商名称', '')
                    if not vendor_name:
                        logging.warning(f"第{i+1}页缺少厂商名称，跳过")
                        continue
                    
                    # 创建厂商文件夹
                    factory_name=clean_factory_name(vendor_name)
                    vendor_folder = make_vendor_folder(factory_name, output_directory)
                    
                    # 提取图片（指定对应的幻灯片编号，i+1对应第几页）
                    img_path = extract_images_from_pptx(file_path, vendor_folder, i+1)
                    if img_path:
                        result['微信'] = img_path
                    else:
                        logging.error(f"提取二维码失败")
                        
                    result['文件路径'] = file_path
                    
                    # 保存结果到厂商文件夹
                    save_result_to_vendor_folder(vendor_folder, result)
                    success_count += 1
                    logging.info(f"第{i+1}页处理成功，厂商：{vendor_name}")
                except Exception as e:
                    logging.error(f"第{i+1}页处理失败: {e}")
            
            if success_count > 0:
                logging.info(f"PPTX文件处理完成，成功处理{success_count}页")
                return True
            else:
                logging.error(f"PPTX文件所有页面处理失败: {file_path}")
                return False
        else:
            logging.error(f"处理文档失败: {file_path}")
            return False
        
    except Exception as e:
        logging.error(f"处理PPTX文件时出错: {e}")
        return False


#---------------------- 程序测试入口 --------------------------------

if __name__ == "__main__":
    # 测试文件路径配置
    file_path = r"tests\ppt\宁波D45期打印资料\温州冠捷科技有限公司.pptx"
    output_directory = r"tests\processed_data\ppt\温州冠捷科技有限公司"
    
    # 执行PPTX处理测试
    if process_pptx_file(file_path,output_directory):
        print("处理成功")
    else:
        print("处理失败")

    
        
