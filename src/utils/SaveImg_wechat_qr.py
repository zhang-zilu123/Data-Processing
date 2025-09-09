import cv2
import os
import uuid
import numpy as np
from pyzbar.pyzbar import decode# 用于识别二维码
import numpy as np
from docx import Document # 用于处理.docx文件
from pptx import Presentation # 用于处理.pptx文件

import fitz # 用于处理pdf文件
from PIL import Image # 用于处理图片

import logging # 用于记录日志

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


#---------------------- 功能性函数--------------------------------

# --- 函数1：将彩色二维码图片转为高对比度黑白二维码图片 ---

def convert_to_black_white_qr(image):
    """
    将彩色二维码图片转为高对比度黑白二维码图片

    """

    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    # Otsu二值化
    _, bw = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    # 判断是否为黑底白码，如果是则反色
    white_ratio = np.sum(bw == 255) / bw.size
    if white_ratio < 0.5:
        bw = 255 - bw
    # 形态学闭运算（可选，防止断裂）
    kernel = np.ones((3, 3), np.uint8)
    bw = cv2.morphologyEx(bw, cv2.MORPH_CLOSE, kernel)
    return bw



# --- 函数2：识别图片是否为黑白微信二维码 ---
def is_wechat_qr_code(image_info):
    """
    识别图片是否为微信二维码

    输入：图片信息 (可以是图片路径字符串或OpenCV图像对象)
    输出：True/False

    """
    # 检查输入是路径还是图像对象
    if isinstance(image_info, str):
        # 如果是路径，则读取图片
        image = cv2.imread(image_info)
    else:
        # 如果已经是OpenCV图像对象 (numpy数组)，则直接使用
        image = image_info

    # 检查图片是否成功读取
    if image is None:
        # print(f"Error: Could not read image.") # 为了处理大量图片时输出更简洁，这里注释掉错误信息
        return False

    # 尝试解码图片中的二维码
    # print(f"尝试黑白解码图片中的二维码")
    decoded_objects = decode(image)

    # 遍历所有解码出的二维码对象
    for obj in decoded_objects:
        # 检查解码出的类型是否为QR码
        if obj.type == 'QRCODE':
            # 将二维码数据解码为UTF-8字符串
            decoded_data = obj.data.decode('utf-8')
            # 简单的判断逻辑：如果解码成功，并且内容可能包含微信相关信息
            # 微信二维码通常会解码出URL或者特定的字符串
            # 这里可以根据实际情况更精确地判断，例如检查URL前缀是否包含“weixin”、“wechat”等
            if "weixin" in decoded_data.lower() or \
            "wechat" in decoded_data.lower() or \
            "mp.weixin.qq.com" in decoded_data.lower() or \
            "wx.qq.com" in decoded_data.lower():
                return True # 识别为微信二维码
        
        
    return False # 未识别为微信二维码





# --- 函数3：水平拼接多张图片 ---
def stitch_images_horizontally(image_list):
    """
    水平拼接多张图片

    输入：list（含多张图片信息，OpenCV图像对象列表）
    输出：一张（拼接）图片信息 (OpenCV图像对象)
    """
    # 如果图片列表为空，则返回None
    if not image_list:
        logging.warning("输入图片列表为空")
        return None

    # 过滤掉列表中可能出现的None值（例如图片读取失败的情况）
    valid_images = [img for img in image_list if img is not None]
    if not valid_images:
        logging.warning("过滤掉列表中可能出现的None值（例如图片读取失败的情况）")
        return None
    try:
        # 找到所有有效图片中高度最大的一张，作为拼接后的最终高度
        max_height = 0
        for img in valid_images:
            max_height = max(max_height, img.shape[0])

        # 调整所有图片的高度以匹配最大高度，并计算拼接后的总宽度
        stitched_width = 0
        resized_images = []
        for img in valid_images:
            # 保持图片原始比例，按最大高度进行缩放
            scale = max_height / img.shape[0]
            resized_img = cv2.resize(img, (int(img.shape[1] * scale), max_height))
            resized_images.append(resized_img)
            stitched_width += resized_img.shape[1] # 累加缩放后的图片宽度

        # 创建一个空白画布，用于拼接（3通道，与彩色图片一致）
        stitched_image = np.zeros((max_height, stitched_width, 3), dtype=np.uint8)

        # 将处理后的图片依次粘贴到空白画布上，实现水平拼接
        current_x = 0
        for img in resized_images:
            stitched_image[0:max_height, current_x:current_x + img.shape[1]] = img
            current_x += img.shape[1]
        
        logging.info(f"图片拼接成功")   
        return stitched_image
    
    except Exception as e:
        logging.error(f"图片拼接过程中发生错误: {str(e)}")
        return None


# --- 函数4：保存图片 ---
def save_image_with_chinese_path(img, save_path):
    ext = os.path.splitext(save_path)[1]
    ret, buf = cv2.imencode(ext, img)
    if ret:
        with open(save_path, 'wb') as f:
            f.write(buf)
        return True
    return False

#----------------------word pptx 图片提取主函数--------------------------------


# --- 函数5：DOCX 图片提取 ---
def extract_images_from_docx(doc_path:str, save_dir:str) -> str:
    """
    从.docx文件中提取图片，保存微信二维码图片，并返回微信二维码图片的存储路径，如果有多张二维码图片，会进行水平拼接后保存。

    处理流程：
    1. 创建保存目录（如果不存在）
    2. 遍历文档中的所有关系，查找图片资源
    3. 将图片数据转换为OpenCV图像格式
    4. 使用二维码识别算法检测微信二维码
    5. 对识别到的二维码进行黑白转换后再次检测
    6. 如果有多张二维码，进行水平拼接处理
    7. 保存最终的二维码图片文件

    
    Args:
        doc_path: docx文件路径
        save_dir (str): 保存图片的目标目录路径
        
    Returns:
        str: 保存的微信二维码图片文件路径，如果没有找到二维码则返回None

    注意事项：
        - 使用UUID生成唯一文件名，避免文件名冲突
        - 支持中文路径的图片保存
        - 自动创建保存目录结构

    """
    if not os.path.exists(save_dir):
        os.makedirs(save_dir)
        # logging.info(f"创建保存目录: {save_dir}")


    wechat_qr_paths = []  # 存储微信二维码图片路径
    wechat_qr_images = []  # 存储微信二维码图片对象
    image_count = 0  # 图片计数器
    qr_count = 0  # 二维码计数器
    save_path = None

    try:

        document = Document(doc_path)
        
        
        # 遍历文档中的所有关系
        for rel in document.part.rels:
            if "image" in document.part.rels[rel].target_ref:
                try:
                    # 获取图片数据
                    image_part = document.part.rels[rel].target_part
                    image_bytes = image_part.blob
                    
                    # 转换为OpenCV图像
                    nparr = np.frombuffer(image_bytes, np.uint8)
                    img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                    
                    if img is not None:
                        image_count += 1
                        
                        
                        # 检查是否为微信二维码
                        if is_wechat_qr_code(img):
                            qr_count += 1
                            wechat_qr_images.append(img)  # 保存图片对象
                        else:
                            bw_img = convert_to_black_white_qr(img)
                            if is_wechat_qr_code(bw_img):
                                qr_count += 1
                                wechat_qr_images.append(img)  # 保存图片对象

                            # logging.info(f"识别到第 {qr_count} 张微信二维码")
                            
                except Exception as e:
                    logging.error(f"处理第 {image_count} 张图片时出错: {str(e)}", exc_info=True)
                    continue
        
        # 处理识别到的二维码图片
        
        if wechat_qr_images:
            
            if len(wechat_qr_images) > 1:
                # 多张二维码图片，进行拼接
                logging.info(f"检测到 {len(wechat_qr_images)} 张二维码图片，开始拼接")
                stitched_image = stitch_images_horizontally(wechat_qr_images)
                if stitched_image is not None:
                    # 保存拼接后的图片
                    
                    filename = f"wechat_qr_{uuid.uuid4().hex}.png"
                    save_path = os.path.join(save_dir, filename)
                    # cv2.imwrite(save_path, stitched_image)
                    success = save_image_with_chinese_path(stitched_image, save_path)
                    
                    if not success:
                        logging.error(f"图片保存失败，路径：{save_path}")
                    else:
                        wechat_qr_paths.append(save_path)
                        logging.info(f"拼接后的二维码图片已保存: {save_path}")

                    
                    
            else:
                # 单张二维码图片，直接保存
                filename = f"wechat_qr_{uuid.uuid4().hex}.png"
                save_path = os.path.join(save_dir, filename)
                success = save_image_with_chinese_path(img, save_path) 
                if not success:
                    logging.error(f"图片保存失败，路径：{save_path}")
                else:
                    logging.info(f"单张二维码图片已保存: {save_path}")

        return save_path
    
    except Exception as e:
        logging.error(f"提取图片过程中发生错误: {str(e)}", exc_info=True)
        
        
        

# --- 函数6：PPTX 图片提取 ---
def extract_images_from_pptx(path:str, save_dir:str, slide_number:int = None) -> str:
    """
    从.pptx文件中提取图片，保存微信二维码图片，并返回微信二维码图片的存储路径。
    如果有多张二维码图片，会进行水平拼接后保存。
    
    处理流程：
    1. 创建保存目录（如果不存在）
    2. 遍历演示文稿中的所有幻灯片
    3. 遍历每张幻灯片中的所有形状，查找包含图片的形状
    4. 将图片数据转换为OpenCV图像格式
    5. 使用二维码识别算法检测微信二维码
    6. 如果有多张二维码，进行水平拼接处理
    7. 保存最终的二维码图片文件

    Args:
        path: pptx文件路径
        save_dir (str): 保存图片的目标目录路径
        
    Returns:
        str: 保存的微信二维码图片文件路径，如果没有找到二维码则返回None

    注意事项：
        - 使用UUID生成唯一文件名，避免文件名冲突
        - 支持多张二维码的水平拼接
        - 自动创建保存目录结构

    """
    
    if not os.path.exists(save_dir):
        os.makedirs(save_dir)
        logging.info(f"创建保存目录: {save_dir}")

    wechat_qr_images = []  # 存储微信二维码图片对象
    wechat_qr_paths = []
    image_count = 0  # 图片计数器
    qr_count = 0  # 二维码计数器
    save_path = None  # 初始化保存路径
    try:
        presentation = Presentation(path)
        
        
        # 遍历演示文稿中的所有幻灯片（或指定幻灯片）
        slides_to_process = presentation.slides
        
        # 如果指定了幻灯片编号，只处理该幻灯片
        if slide_number is not None:
            if 1 <= slide_number <= len(presentation.slides):
                slides_to_process = [presentation.slides[slide_number - 1]]  # 转换为0-based索引
                logging.info(f"只处理第 {slide_number} 张幻灯片")
            else:
                logging.error(f"指定的幻灯片编号 {slide_number} 超出范围 (1-{len(presentation.slides)})")
                return None
        else:
            logging.info(f"处理所有 {len(presentation.slides)} 张幻灯片")
        
        for slide_num, slide in enumerate(slides_to_process, 1 if slide_number is None else slide_number):
            logging.debug(f"处理第 {slide_num} 张幻灯片")
            
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                # 如果形状包含图像属性
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        # 将二进制数据转换为numpy数组，再用cv2解码成OpenCV图像
                        nparr = np.frombuffer(image_bytes, np.uint8)
                        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                        
                        if img is not None:
                            image_count += 1
                            
                            # 检查是否为微信二维码
                            if is_wechat_qr_code(img):
                                qr_count += 1
                                wechat_qr_images.append(img)
                            else:
                                bw_img = convert_to_black_white_qr(img)
                                if is_wechat_qr_code(bw_img):
                                    qr_count += 1
                                    wechat_qr_images.append(img)

                                
                    except Exception as e:
                        logging.error(f"处理第 {image_count} 张图片时出错: {str(e)}", exc_info=True)
                        continue

        # 处理识别到的二维码图片
        
        if wechat_qr_images:
            
            if len(wechat_qr_images) > 1:
                # 多张二维码图片，进行拼接
                logging.info(f"检测到 {len(wechat_qr_images)} 张二维码图片，开始拼接")
                stitched_image = stitch_images_horizontally(wechat_qr_images)
                if stitched_image is not None:
                    # 保存拼接后的图片
                    
                    filename = f"wechat_qr_{uuid.uuid4().hex}.png"
                    save_path = os.path.join(save_dir, filename)
                    success = save_image_with_chinese_path(stitched_image, save_path)
                    
                    if not success:
                        logging.error(f"图片保存失败，路径：{save_path}")
                    else:
                        wechat_qr_paths.append(save_path)
                        logging.info(f"拼接后的二维码图片已保存: {save_path}")

                    
                    
            else:
                # 单张二维码图片，直接保存
                filename = f"wechat_qr_{uuid.uuid4().hex}.png"
                save_path = os.path.join(save_dir, filename)
                success = save_image_with_chinese_path(img, save_path)
                    
                if not success:
                    logging.error(f"图片保存失败，路径：{save_path}")
                else:
                    logging.info(f"单张二维码图片已保存: {save_path}")

        
        logging.info(f"从PPTX文稿中提取到 {image_count} 张图片，其中 {qr_count} 张是微信二维码")
        return save_path
    
    except Exception as e:
        logging.error(f"提取图片过程中发生错误: {str(e)}", exc_info=True)
        

# --- 函数7：PDF 图片提取 ---
def extract_images_from_pdf(file_path: str, output_dir: str, page_num: int) -> str:
    """

    从PDF文件中提取指定页的图片，识别并保存微信二维码图片。
    如果找到多个二维码图片，会返回第一个识别的二维码图片路径。
    
    处理流程：
    1. 打开PDF文件并加载指定页面
    2. 提取页面中的所有图片
    3. 将每张图片临时保存到磁盘
    4. 使用二维码识别算法检测微信二维码
    5. 如果是二维码则保留并重命名，非二维码图片则删除
    6. 返回第一个找到的二维码图片路径

    参数:
        file_path (str): PDF文件路径
        output_dir (str): 保存图片的目标目录路径
        page_num (int): 要提取的页码(1-based)
        
    返回:
        str: 保存的微信二维码图片文件路径，如果没有找到二维码则返回空字符串

    注意事项：
        - 使用UUID生成唯一文件名，避免文件名冲突
        - 页码是1-based，但内部处理会转换为0-based
        - 仅返回第一个识别的二维码图片
        - 会自动清理非二维码的临时图片文件

    """
    
    # 添加这行：确保输出目录存在
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        logging.info(f"创建保存目录: {output_dir}")
    
    
    try:
        pdf_doc = fitz.open(file_path)
        wechat_qr_images = []
        qr_count = 0
        image_count = 0
        save_path = None
        
        # 加载指定页面
        page = pdf_doc.load_page(page_num - 1)  # 0-based
        images = page.get_images()
        
        if not images:
            logging.warning(f"第{page_num}页没有图片")
            return ""
            
        for img_index, img in enumerate(images):
            try:
                xref = img[0]
                base_image = pdf_doc.extract_image(xref)
                img_data = base_image["image"]
                
                # 转换为OpenCV图像
                nparr = np.frombuffer(img_data, np.uint8)
                cv_img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                
                if cv_img is not None:
                    image_count += 1
                    
                    # 检查是否为微信二维码
                    if is_wechat_qr_code(cv_img):
                        qr_count += 1
                        wechat_qr_images.append(cv_img)  # 保存图片对象
                    else:
                        bw_img = convert_to_black_white_qr(cv_img)
                        if is_wechat_qr_code(bw_img):
                            qr_count += 1
                            wechat_qr_images.append(cv_img)  # 保存图片对象

                        # logging.info(f"识别到第 {qr_count} 张微信二维码")
                        
            except Exception as e:
                logging.error(f"处理第 {image_count} 张图片时出错: {str(e)}", exc_info=True)
                continue
            # 处理识别到的二维码图片
        if wechat_qr_images:
            if len(wechat_qr_images) > 1:
                # 多张二维码图片，进行拼接
                logging.info(f"检测到 {len(wechat_qr_images)} 张二维码图片，开始拼接")
                stitched_image = stitch_images_horizontally(wechat_qr_images)
                if stitched_image is not None:
                    # 保存拼接后的图片
                    filename = f"wechat_qr_{uuid.uuid4().hex}.png"
                    save_path = os.path.join(output_dir, filename)
                    success = save_image_with_chinese_path(stitched_image, save_path)
                    
                    if not success:
                        logging.error(f"图片保存失败，路径：{save_path}")
                    else:
                        logging.info(f"拼接后的二维码图片已保存: {save_path}")
            else:
                # 单张二维码图片，直接保存
                filename = f"wechat_qr_{uuid.uuid4().hex}.png"
                save_path = os.path.join(output_dir, filename)
                success = save_image_with_chinese_path(wechat_qr_images[0], save_path)
                
                if not success:
                    logging.error(f"图片保存失败，路径：{save_path}")
                else:
                    logging.info(f"单张二维码图片已保存: {save_path}")
        
        pdf_doc.close()
        logging.info(f"从PDF第{page_num}页提取到 {image_count} 张图片，其中 {qr_count} 张是微信二维码")
        return save_path if save_path else ""
    
    except Exception as e:
        logging.error(f"提取PDF图片过程中发生错误: {str(e)}", exc_info=True)
        return ""
        

            
    # page = doc.load_page(page_num - 1)  # 0-based
    # images = page.get_images()
    
    # if not images:
    #     logging.warning(f"没有图片")
    #     return ""
        
    # for img_index, img in enumerate(images):
    #     xref = img[0]
    #     base_image = doc.extract_image(xref)
    #     img_data = base_image["image"]
        
        
        
    #     # 保存临时图片
    #     img_path = os.path.join(output_dir, f"temp_img_{page_num}_{img_index}.png")
    #     if os.path.exists(img_path):
    #         print(f"保存成功:{img_path}")


    #     with open(img_path, "wb") as f:  # 这里就不会报错了
    #         f.write(img_data)
        
    #     # 检查是否是二维码
    #     try:
    #         pil_img = Image.open(img_path)
    #         # decoded_qr = decode(pil_img)
    #         decoded_qr = is_wechat_qr_code(pil_img)


    #         print(f"decoded_qr: {decoded_qr}")
    #         if decoded_qr:  # 是二维码
    #             final_path = os.path.join(output_dir, f"wechat_qr_{uuid.uuid4().hex}.png")
    #             os.rename(img_path, final_path)
    #             return final_path
    #         else:  # 不是二维码，删除临时文件
    #             print("不是二维码")
    #             # os.remove(img_path)
    #     except Exception as e:
    #         logging.error(f"二维码解析失败: {e}")
    #         if os.path.exists(img_path):
    #             os.remove(img_path)
    
    return ""


# --- 示例用法 ---
if __name__ == "__main__":

    # 测试1：从docx提取图片
    
    # test_docx_path = r"tests\word\永康市聚野工贸有限公司.docx"
    # output_dir = r"tests\word\永康市聚野工贸有限公司"
    # print("测试1：从docx提取图片")
    # try:
       
    #     images_from_docx = extract_images_from_docx(test_docx_path, output_dir)
    #     if images_from_docx:
    #         print("是微信二维码")
    #     else:
    #         print("不是微信二维码")
        
    # except Exception as e:
    #     logging.error(f"从docx提取图片失败: {str(e)}", exc_info=True)
    
    
    # 测试2：从pptx提取图片

    test_pptx_path = r"data\input_data\ppt\广东鸿祺玩具实业有限公司.pptx"
    output_dir = r"tests\processed_data\ppt\广东鸿祺玩具实业有限公司"
    
    
    qr_paths = extract_images_from_pptx(test_pptx_path, output_dir, slide_number=1)

    # 测试3：从pdf提取图片
    
    # test_pptx_path = r"data\input_data\ppt\东莞市安琦家居用品有限公司.pptx"
    # output_dir = r"tests\processed_data\pdf\青岛柔百纳纸制品有限公司1750299931"
    
    
    # qr_paths = extract_images_from_pdf(test_pptx_path, output_dir, 1)